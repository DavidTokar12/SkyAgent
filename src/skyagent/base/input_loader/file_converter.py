from __future__ import annotations

import base64
import csv
import logging

from io import StringIO
from typing import TYPE_CHECKING

import docx
import fitz
import pandas as pd

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from skyagent.base.exceptions import SkyAgentFileError


if TYPE_CHECKING:
    from pathlib import Path


logger = logging.getLogger(__name__)


def default_doc_converter(file_path: Path) -> tuple[str, list[str]]:
    """
    Extract text, tables, and images from a Word document.

    Args:
        file_path: Path to the Word file

    Returns:
        Tuple containing:
            - Extracted text content including tables in CSV format
            - List of base64 encoded images
    """
    doc = docx.Document(file_path)
    text_parts = []
    images = []

    try:
        for element in doc.element.body:

            if element.tag.endswith("p"):
                if len(element.text.strip()) > 0:
                    text_parts.append(element.text)
            elif element.tag.endswith("tbl"):
                output = StringIO()
                writer = csv.writer(output)

                table_index = sum(
                    1 for e in element.getprevious() if e.tag.endswith("tbl")
                )
                table = doc.tables[table_index]

                for row in table.rows:
                    writer.writerow(cell.text for cell in row.cells)

                text_parts.append(output.getvalue().strip())

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_bytes = rel.target_part.blob
                    base64_image = base64.b64encode(image_bytes).decode("utf-8")
                    images.append(base64_image)
                except Exception as e:
                    logger.warning(f"Failed to extract image {rel.target_ref}: {e}")
                    continue

        return ["\n".join(text_parts)], images

    except Exception as e:
        raise SkyAgentFileError(f"Failed to process Word file: {e}")


def default_csv_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Convert spreadsheet files (CSV, XLS, XLSX) to list of CSV contents.

    Args:
        file_path: Path to the spreadsheet file

    Returns:
        Tuple containing:
            - List of CSV contents (one per sheet)
            - Empty list (spreadsheets don't contain images)
    """
    try:
        suffix = file_path.suffix.lower()

        if suffix == ".csv":
            df = pd.read_csv(file_path)
            return [df.to_csv(index=False)], []
        else:
            sheets_dict = pd.read_excel(file_path, sheet_name=None)
            return [df.to_csv(index=False) for df in sheets_dict.values()], []

    except Exception as e:
        raise SkyAgentFileError(f"Failed to process spreadsheet file: {e}")


def default_pdf_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Extract text and images from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        Tuple containing:
            - List containing a single string with all text content
            - List of base64 encoded images
    """
    text_content = []
    images = []

    try:
        pdf_document = fitz.open(file_path)

        for page in pdf_document:

            # TODO also extract tables correctly
            text_content.append(page.get_text())

            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]

                    base64_image = base64.b64encode(image_bytes).decode("utf-8")
                    images.append(base64_image)

                except Exception as e:
                    logger.warning(
                        f"Failed to extract image {img_index} from page {page.number + 1}: {e}"
                    )
                    continue

        return ["\n".join(text_content)], images

    except Exception as e:
        raise SkyAgentFileError(f"Failed to process PDF file: {e}")
    finally:
        if "pdf_document" in locals():
            pdf_document.close()


def default_ppt_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Extract text, tables, and images from a PowerPoint presentation.
    Combines all slides into a single text content.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        Tuple containing:
            - List containing a single string with all slides' text and tables
            - List of base64 encoded images
    """
    text_parts = []
    images = []

    try:
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {slide_num}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    output = StringIO()
                    writer = csv.writer(output)

                    for row in shape.table.rows:
                        writer.writerow(cell.text for cell in row.cells)

                    text_parts.append(output.getvalue().strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        base64_image = base64.b64encode(image_bytes).decode("utf-8")
                        images.append(base64_image)
                    except Exception as e:
                        logger.warning(
                            f"Failed to extract image from slide {slide_num}: {e}"
                        )
                        continue

        return ["\n\n".join(text_parts)], images

    except Exception as e:
        raise SkyAgentFileError(f"Failed to process PowerPoint file: {e}")


def default_img_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Convert image file to base64 string.

    Args:
        file_path: Path to the image file (jpg, jpeg, or png)

    Returns:
        Tuple containing:
            - List with single base64 string representation of the image
            - Empty list (we don't extract images from images)
    """
    try:
        with open(file_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
            return [encoded_string], []

    except Exception as e:
        raise SkyAgentFileError(f"Failed to process image file: {e}")
