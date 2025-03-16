from __future__ import annotations

import base64
import csv
import logging
import os
import re
import urllib.error
import urllib.request

from io import BytesIO
from io import StringIO
from pathlib import Path

import fitz
import html2text
import mammoth
import pandas as pd

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from skyagent.base.exceptions import SkyAgentFileError


logger = logging.getLogger(__name__)


def _is_url_image(image_src: str) -> bool:
    """Check if the image source is a URL."""
    return image_src.startswith(("http://", "https://"))


def _is_base64_image(image_src: str) -> bool:
    """Check if the image source is a base64 encoded image."""
    return image_src.startswith("data:image")


def _read_image_to_base64(image_path: Path) -> str:
    """
    Read an image file from disk and return it as a base64-encoded JPEG (re-encoded to JPEG).
    """
    try:
        with Image.open(image_path) as im:
            rgb_im = im.convert("RGB")
            output = BytesIO()
            rgb_im.save(output, format="JPEG")
            jpeg_data = output.getvalue()

        return f'data:image/jpeg;base64,{base64.b64encode(jpeg_data).decode("utf-8")}'
    except Exception as e:
        raise SkyAgentFileError("Failed to read image file") from e


def _convert_to_jpeg_base64(base64_str: str) -> str:
    """
    Convert any base64-encoded image to JPEG format and return as base64.
    Always re-encodes the image to JPEG.
    """
    if "," in base64_str:
        _, data = base64_str.split(",", 1)
    else:
        raise SkyAgentFileError("Failed to convert image to JPEG")

    try:
        image_data = base64.b64decode(data)
        with Image.open(BytesIO(image_data)) as img:
            rgb_img = img.convert("RGB")

            output = BytesIO()
            rgb_img.save(output, format="JPEG")
            jpeg_data = output.getvalue()

        return f'data:image/jpeg;base64,{base64.b64encode(jpeg_data).decode("utf-8")}'
    except Exception as e:
        raise SkyAgentFileError("Failed to convert image to JPEG") from e


def _download_image_to_base64(url: str) -> str:
    """
    Download an image from a URL and return it as a base64-encoded JPEG (re-encoded to JPEG).
    """
    try:
        with urllib.request.urlopen(url) as response:
            image_data = response.read()

        with Image.open(BytesIO(image_data)) as im:
            rgb_im = im.convert("RGB")
            output = BytesIO()
            rgb_im.save(output, format="JPEG")
            jpeg_data = output.getvalue()

        return f'data:image/jpeg;base64,{base64.b64encode(jpeg_data).decode("utf-8")}'
    except (urllib.error.URLError, Exception) as e:
        raise SkyAgentFileError("Failed to download image from URL") from e


def _extract_images_from_markdown(
    markdown_content: str, original_file_path: str | Path
) -> tuple[str, list[str]]:
    """
    Extract images from markdown content and return cleaned markdown and base64 encoded JPEGs.

    Args:
        markdown_content (str): The original markdown content
        original_file_path (str | Path): The path to the original markdown file

    Returns:
        tuple[str, list[str]]: (cleaned markdown content, list of base64 encoded JPEG images)
    """
    original_path = Path(original_file_path)
    base_path = original_path.parent

    image_pattern = r"!\[([^\]]*)\]\(([^)]+)\)"
    images = []

    def handle_images(match):
        alt_text, image_src = match.groups()

        if _is_base64_image(image_src):

            try:
                jpeg_base64 = _convert_to_jpeg_base64(image_src)
                images.append(jpeg_base64)
            except SkyAgentFileError as e:
                logger.error(f"Error converting base64 image: {e}")

            return f'![{alt_text}]({image_src.split(",")[0]},)'

        elif _is_url_image(image_src):
            try:
                base64_data = _download_image_to_base64(image_src)
                images.append(base64_data)
            except SkyAgentFileError as e:
                logger.error(f"Error downloading image from URL {image_src}: {e}")

            return match.group(0)  # Keep original path in markdown

        else:
            if os.path.isabs(image_src):
                image_path = Path(image_src)
            else:
                image_path = base_path / image_src

            if image_path.exists():
                try:
                    base64_data = _read_image_to_base64(image_path)
                    images.append(base64_data)
                except SkyAgentFileError as e:
                    logger.error(f"Error reading image file {image_path}: {e}")
            else:
                logger.error(
                    f"Could not find image file: {image_path}. Image will not be included in output."
                )
            return match.group(0)  # Keep original path in markdown

    cleaned_markdown = re.sub(image_pattern, handle_images, markdown_content)
    return cleaned_markdown, images


def _extract_text_and_images_from_html(
    html_content: str, original_file_path: str | Path
) -> tuple[str, list[str]]:
    """
    Convert the given HTML to Markdown using html2text, extracting images
    as base64 for three cases:
      1) data:image/... (already base64-encoded in the HTML)
      2) local file path (read from disk and encode)
      3) remote URL (download and encode)

    Args:
        html_content (str): The HTML string to process.
        original_file_path (str|Path|None): Path to the original file.

    Returns:
        (md_text, base64_images):
            md_text (str): The HTML converted to Markdown (via html2text).
            base64_images (List[str]): A list of base64-encoded image strings.
    """
    try:
        text_maker = html2text.HTML2Text()
        text_maker.ignore_links = False
        text_maker.bypass_tables = False
        text_maker.ignore_images = False

        md_text = text_maker.handle(html_content)

        cleaned_markdown, base64_images = _extract_images_from_markdown(
            md_text, original_file_path
        )

        return [cleaned_markdown], base64_images
    except Exception as e:
        raise SkyAgentFileError("Failed to extract text and images from HTML.") from e


def default_doc_converter(file_path: Path) -> tuple[str, list[str]]:
    """
    Extract text, tables, and images from a Word document.

    Args:
        file_path: Path to the Word file

    Returns:
        Tuple containing:
            - Extracted text content (as a single string)
            - List of base64 encoded images
    """
    try:
        with open(file_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)

        html_content = result.value

        return _extract_text_and_images_from_html(
            html_content=html_content, original_file_path=file_path
        )
    except FileNotFoundError:
        raise SkyAgentFileError(f"Word document not found: {file_path}")
    except SkyAgentFileError:
        # Re-raise SkyAgentFileError from _extract_text_and_images_from_html
        raise
    except Exception as e:
        raise SkyAgentFileError("Failed to convert Word document.") from e


def default_xls_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Convert (XLS, XLSX) to list of CSV contents.

    Args:
        file_path: Path to the spreadsheet file

    Returns:
        Tuple containing:
            - List of CSV contents (one per sheet)
            - Empty list (spreadsheets don't contain images)
    """
    try:
        sheets_dict = pd.read_excel(file_path, sheet_name=None)
        csv_sheets = []

        for sheet_name, df in sheets_dict.items():

            try:
                csv_sheets.append(df.to_csv(index=False))
            except Exception as e:
                logger.error(f"Error converting sheet '{sheet_name}' to CSV: {e!s}")

        return csv_sheets, []

    except FileNotFoundError:
        raise SkyAgentFileError(f"Spreadsheet file not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError("Failed to process spreadsheet file.") from e


def default_pdf_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Extract text and images from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        Tuple containing:
            - List containing a single string with all text content
            - List of base64 encoded images
    """
    text_content = []
    images = []

    try:
        with fitz.open(file_path) as pdf_document:
            for page_idx, page in enumerate(pdf_document):
                try:
                    page_text = page.get_text()
                    text_content.append(page_text)
                except Exception as e:
                    logger.error(
                        f"Failed to extract text from page {page_idx+1}: {e!s}"
                    )

                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = pdf_document.extract_image(xref)
                        image_bytes = base_image["image"]

                        base64_image = base64.b64encode(image_bytes).decode("utf-8")

                        images.append(base64_image)
                    except Exception as e:
                        logger.warning(
                            f"Failed to extract image {img_index+1} from page {page_idx+1}: {e!s}"
                        )
                        continue

        return ["\n".join(text_content)], images

    except FileNotFoundError:
        raise SkyAgentFileError(f"PDF file not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError("Failed to process PDF file.") from e


def default_ppt_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Extract text, tables, and images from a PowerPoint presentation.
    Combines all slides into a single text content.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        Tuple containing:
            - List containing a single string with all slides' text and tables
            - List of base64 encoded images
    """
    text_parts = []
    images = []

    try:
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                try:

                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        output = StringIO()
                        writer = csv.writer(output)

                        for row in shape.table.rows:
                            writer.writerow(cell.text for cell in row.cells)

                        text_parts.append(output.getvalue().strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_bytes = image.blob
                        base64_image = base64.b64encode(image_bytes).decode("utf-8")
                        images.append(base64_image)

                except Exception as e:
                    logger.warning(
                        f"Error processing shape on slide {slide_num}: {e!s}"
                    )
                    continue

        return text_parts, images

    except FileNotFoundError:
        raise SkyAgentFileError(f"PowerPoint file not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError("Failed to process PowerPoint file.") from e


def default_img_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Convert image file to base64 string.

    Args:
        file_path: Path to the image file (jpg, jpeg, or png)

    Returns:
        Tuple containing:
            - List with single base64 string representation of the image
            - Empty list (we don't extract images from images)
    """
    try:
        with open(file_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
            return [encoded_string], []
    except FileNotFoundError:
        raise SkyAgentFileError(f"Image file not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError("Failed to process image file.") from e


def _no_special_conversion(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Generic text reader: reads file content as plain text.

    Returns:
        (list_of_text, empty_list_for_images)
    """
    try:
        file_content = file_path.read_text()
        return [file_content], []
    except FileNotFoundError:
        raise SkyAgentFileError(f"File not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError("Failed to read text from file.") from e


def default_code_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)


def default_markdown_converter(file_path: Path) -> tuple[list[str], list[str]]:
    """
    Extract images from markdown content (base64, URLs, or local files),
    and return the cleaned markdown plus a list of extracted base64 images.
    """
    try:
        markdown_content = file_path.read_text()
        cleaned_markdown, base64_images = _extract_images_from_markdown(
            markdown_content=markdown_content, original_file_path=file_path
        )
        return [cleaned_markdown], base64_images
    except FileNotFoundError:
        raise SkyAgentFileError(f"Markdown file not found: {file_path}")
    except Exception as e:
        raise SkyAgentFileError(f"Failed to process markdown file: {e}") from e


def default_json_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)


def default_yaml_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)


def default_xml_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)


def default_text_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)


def default_csv_converter(file_path: Path) -> tuple[list[str], list[str]]:
    return _no_special_conversion(file_path=file_path)
